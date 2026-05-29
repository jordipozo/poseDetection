"""Generate PowerPoint presentation for PoseDetection app."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as pns
from lxml import etree

# ── Palette ─────────────────────────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x0D, 0x0D)   # #0d0d0d  main bg
BG_CARD   = RGBColor(0x1A, 0x1A, 0x1A)   # #1a1a1a  surface
RED       = RGBColor(0xE5, 0x39, 0x35)   # #e53935  accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CYAN      = RGBColor(0x00, 0xBC, 0xD4)   # face landmarks color
GREY_LT   = RGBColor(0xBB, 0xBB, 0xBB)
GREY_MED  = RGBColor(0x77, 0x77, 0x77)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helpers ──────────────────────────────────────────────────────────────────

def fill_bg(slide, color: RGBColor):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_SHAPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=Pt(13), title=None, title_color=RED,
                   bullet="▸ ", text_color=GREY_LT):
    """Add a list of bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = font_size + Pt(2)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet + item
        run.font.size = font_size
        run.font.color.rgb = text_color
    return txBox


def accent_bar(slide, color=RED, height=Inches(0.07)):
    """Thin coloured bar at top of slide."""
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=color)


def slide_number(slide, n, total):
    add_text(slide, f"{n} / {total}",
             SLIDE_W - Inches(1.1), SLIDE_H - Inches(0.4),
             Inches(1), Inches(0.35),
             font_size=Pt(10), color=GREY_MED, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank_layout)
fill_bg(s1, BG_DARK)
accent_bar(s1)

# decorative vertical stripe
add_rect(s1, Inches(0.5), Inches(1.2), Inches(0.06), Inches(4.8), fill_color=RED)

# title
add_text(s1, "PoseDetection",
         Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
         font_size=Pt(54), bold=True, color=WHITE)

# subtitle
add_text(s1, "Captura y Análisis de Movimiento en Baloncesto",
         Inches(0.8), Inches(3.0), Inches(11), Inches(0.8),
         font_size=Pt(22), color=CYAN)

# description
add_text(s1,
         "Aplicación web progresiva para investigación doctoral en biomecánica deportiva.\n"
         "Detección de pose humana en tiempo real · 33 landmarks · Exportación de metadatos JSON.",
         Inches(0.8), Inches(3.9), Inches(10), Inches(1.4),
         font_size=Pt(14), color=GREY_LT)

# bottom info
add_text(s1, "Doctorado UVa  ·  2026", Inches(0.8), Inches(6.8), Inches(6), Inches(0.5),
         font_size=Pt(11), color=GREY_MED)

slide_number(s1, 1, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Contexto y Motivación
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank_layout)
fill_bg(s2, BG_DARK)
accent_bar(s2)

add_text(s2, "Contexto y Motivación", Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# left card
add_rect(s2, Inches(0.4), Inches(1.0), Inches(5.8), Inches(5.6), fill_color=BG_CARD)
add_bullet_box(s2,
    ["Análisis biomecánico de jugadores de baloncesto",
     "Necesidad de datos de movimiento a nivel de fotograma",
     "Ausencia de herramientas offline accesibles en campo",
     "Investigación doctoral sobre rendimiento deportivo"],
    Inches(0.55), Inches(1.1), Inches(5.5), Inches(2.8),
    title="Problema", font_size=Pt(13))

add_bullet_box(s2,
    ["Captura de vídeo + esqueleto en tiempo real",
     "Metadatos JSON completos por fotograma",
     "Funciona en móvil sin instalación (PWA)",
     "Etiquetado de movimientos de baloncesto (7 tipos)"],
    Inches(0.55), Inches(3.6), Inches(5.5), Inches(2.6),
    title="Solución", font_size=Pt(13))

# right card
add_rect(s2, Inches(6.6), Inches(1.0), Inches(6.3), Inches(5.6), fill_color=BG_CARD)
add_text(s2, "Movimientos Capturados",
         Inches(6.75), Inches(1.15), Inches(6.0), Inches(0.45),
         font_size=Pt(15), bold=True, color=RED)

movements = [
    ("1", "Tiro libre"),
    ("2", "Triple"),
    ("3", "Salto"),
    ("4", "Rebote"),
    ("5", "Defensa"),
    ("6", "Cambio de dirección"),
    ("7", "Bandeja"),
]
for i, (code, name) in enumerate(movements):
    y = Inches(1.7) + i * Inches(0.64)
    add_rect(s2, Inches(6.75), y, Inches(0.42), Inches(0.44), fill_color=RED)
    add_text(s2, code, Inches(6.76), y + Pt(4), Inches(0.42), Inches(0.44),
             font_size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s2, name, Inches(7.3), y + Pt(4), Inches(5.3), Inches(0.44),
             font_size=Pt(13), color=GREY_LT)

slide_number(s2, 2, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Arquitectura General
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank_layout)
fill_bg(s3, BG_DARK)
accent_bar(s3)

add_text(s3, "Arquitectura de la Aplicación", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Architecture: pipeline boxes
pipeline = [
    ("Cámara\n(getUserMedia)", CYAN),
    ("<video>\noculto", GREY_MED),
    ("RAF\nLoop", RED),
    ("MediaPipe\nPose (WASM)", RED),
    ("Canvas\n2D Render", CYAN),
    ("MediaRecorder\nVideo", RGBColor(0xFF,0xA0,0x00)),
    ("Ficheros\n(Video + JSON)", RGBColor(0x4C,0xAF,0x50)),
]

box_w = Inches(1.55)
box_h = Inches(1.0)
gap   = Inches(0.22)
total_w = len(pipeline) * box_w + (len(pipeline)-1) * gap
start_x = (SLIDE_W - total_w) / 2
y_box = Inches(1.4)

for i, (label, color) in enumerate(pipeline):
    x = start_x + i * (box_w + gap)
    add_rect(s3, x, y_box, box_w, box_h, fill_color=color)
    add_text(s3, label, x, y_box + Pt(6), box_w, box_h,
             font_size=Pt(11), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        ax = x + box_w + Pt(3)
        add_text(s3, "→", ax, y_box + Pt(20), gap, box_h,
                 font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Three component cards below
cards = [
    ("Motor de Detección",
     ["Google MediaPipe Pose",
      "33 landmarks por persona",
      "Inferencia: WebAssembly + WebGL",
      "Sin servidor — 100% en navegador",
      "Umbral visibilidad: 0.5"]),
    ("Pipeline de Grabación",
     ["canvas.captureStream() → MediaRecorder",
      "Chunks cada 200 ms",
      "MP4 (H.264) / WebM (VP9/VP8)",
      "Resolución fija: 640 × 480 px",
      "Límite: 30 segundos"]),
    ("Gestión de Ficheros",
     ["File System Access API (Chromium)",
      "Auto-descarga (desktop)",
      "Descarga manual (iOS/Firefox)",
      "Naming: {dorsal}_{ts}_{postura}",
      "JSON con 6 decimales de precisión"]),
]

card_w = Inches(3.9)
card_h = Inches(3.55)
card_y = Inches(2.65)
card_gap = Inches(0.37)
cx = (SLIDE_W - 3*card_w - 2*card_gap) / 2

for i, (title, items) in enumerate(cards):
    x = cx + i * (card_w + card_gap)
    add_rect(s3, x, card_y, card_w, card_h, fill_color=BG_CARD)
    add_bullet_box(s3, items, x + Inches(0.12), card_y + Inches(0.05),
                   card_w - Inches(0.2), card_h - Inches(0.1),
                   title=title, font_size=Pt(12))

slide_number(s3, 3, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Stack Tecnológico
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank_layout)
fill_bg(s4, BG_DARK)
accent_bar(s4)

add_text(s4, "Stack Tecnológico", Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

tech_groups = [
    ("Núcleo", WHITE, [
        "HTML5 · CSS3 · JavaScript ES2020+",
        "Single Page Application sin build tools",
        "Un único fichero: index.html (~1200 líneas)",
    ]),
    ("Modelo IA", CYAN, [
        "MediaPipe Pose (Google) vía CDN jsDelivr",
        "Model complexity: 1 (equilibrado)",
        "Detección 0.5 · Tracking 0.5 (umbral)",
    ]),
    ("APIs del Navegador", RGBColor(0xFF,0xA0,0x00), [
        "MediaDevices.getUserMedia() — cámara",
        "Canvas 2D API — render en tiempo real",
        "MediaRecorder API — codificación vídeo",
        "File System Access API — escritura disco",
        "performance.now() — timestamps precisos",
    ]),
    ("PWA / Offline", RGBColor(0x4C,0xAF,0x50), [
        "Service Worker — cache-first strategy",
        "manifest.json — instalable en dispositivo",
        "Pre-cache libs MediaPipe (~8-15 MB)",
        "Modo fullscreen, portrait-primary",
    ]),
]

col_w = Inches(5.9)
col_h = Inches(5.6)
left_x = Inches(0.4)
right_x = Inches(6.8)
top_y = Inches(1.05)
mid_y = Inches(3.9)

positions = [
    (left_x, top_y), (right_x, top_y),
    (left_x, mid_y), (right_x, mid_y),
]
heights = [Inches(2.55), Inches(2.55), Inches(2.95), Inches(2.95)]

for (x, y), h, (title, color, items) in zip(positions, heights, tech_groups):
    add_rect(s4, x, y, col_w, h, fill_color=BG_CARD)
    add_bullet_box(s4, items, x + Inches(0.12), y + Inches(0.05),
                   col_w - Inches(0.2), h - Inches(0.1),
                   title=title, title_color=color, font_size=Pt(12))

slide_number(s4, 4, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – MediaPipe en Detalle
# ═══════════════════════════════════════════════════════════════════════════════
s5b = prs.slides.add_slide(blank_layout)
fill_bg(s5b, BG_DARK)
accent_bar(s5b)

add_text(s5b, "MediaPipe Pose — Cómo Funciona", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Header strip: what is MediaPipe
add_rect(s5b, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.1), fill_color=BG_CARD)
add_text(s5b,
         "MediaPipe es un framework de ML de Google que ejecuta modelos de visión artificial directamente "
         "en el dispositivo, sin enviar datos a ningún servidor. En PoseDetection se usa el modelo "
         "BlazePose, diseñado para detectar y seguir 33 puntos clave del cuerpo humano en tiempo real.",
         Inches(0.55), Inches(1.05), Inches(12.1), Inches(1.0),
         font_size=Pt(13), color=GREY_LT, wrap=True)

# Three columns: Model, Execution, Output
col3_w = Inches(3.9)
col3_h = Inches(4.4)
col3_y = Inches(2.3)
col3_gap = Inches(0.37)
col3_x = (SLIDE_W - 3*col3_w - 2*col3_gap) / 2

cols_data = [
    ("BlazePose — El Modelo", CYAN, [
        "Red neuronal CNN de dos etapas",
        "Etapa 1: detector de persona (bounding box)",
        "Etapa 2: regresión de 33 landmarks",
        "Entrenado con miles de vídeos deportivos",
        "Model complexity 1: equilibrio velocidad/precisión",
        "Alternativas: 0 (rápido) y 2 (preciso)",
    ]),
    ("Ejecución en Navegador", RGBColor(0xFF,0xA0,0x00), [
        "WebAssembly (WASM): código nativo en browser",
        "WebGL: acelera con la GPU del dispositivo",
        "Sin instalación — carga desde CDN jsDelivr",
        "Primera carga: ~8-15 MB de modelo",
        "Caché offline via Service Worker",
        "Llamada: pose.send(videoElement) por frame",
    ]),
    ("Salida del Modelo", RGBColor(0x4C,0xAF,0x50), [
        "Callback onResults() por cada frame",
        "poseLandmarks: coordenadas 2D normalizadas",
        "poseWorldLandmarks: coordenadas 3D reales",
        "x, y: posición relativa al frame [0-1]",
        "z: profundidad relativa a las caderas",
        "visibility: confianza de detección [0-1]",
    ]),
]

for i, (title, color, items) in enumerate(cols_data):
    x = col3_x + i * (col3_w + col3_gap)
    add_rect(s5b, x, col3_y, col3_w, col3_h, fill_color=BG_CARD)
    add_bullet_box(s5b, items, x + Inches(0.12), col3_y + Inches(0.05),
                   col3_w - Inches(0.2), col3_h - Inches(0.1),
                   title=title, title_color=color, font_size=Pt(12))

# Bottom note
add_text(s5b,
         "▸  Umbral en PoseDetection: detection_confidence = 0.5  ·  tracking_confidence = 0.5  "
         "·  landmarks con visibility < 0.3 no se dibujan (pero sí se guardan en el JSON)",
         Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.5),
         font_size=Pt(11), color=GREY_MED, wrap=True)

slide_number(s5b, 5, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – PWA en Detalle
# ═══════════════════════════════════════════════════════════════════════════════
s5c = prs.slides.add_slide(blank_layout)
fill_bg(s5c, BG_DARK)
accent_bar(s5c)

add_text(s5c, "Progressive Web App (PWA)", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Intro strip
add_rect(s5c, Inches(0.4), Inches(1.0), Inches(12.5), Inches(1.1), fill_color=BG_CARD)
add_text(s5c,
         "Una PWA es una aplicación web que adopta capacidades de las apps nativas: funciona offline, "
         "puede instalarse en el dispositivo y se abre en pantalla completa. PoseDetection es una PWA "
         "de un solo fichero HTML — no requiere tienda de apps, backend ni proceso de instalación.",
         Inches(0.55), Inches(1.05), Inches(12.1), Inches(1.0),
         font_size=Pt(13), color=GREY_LT, wrap=True)

# Left: Service Worker
add_rect(s5c, Inches(0.4), Inches(2.3), Inches(6.0), Inches(4.4), fill_color=BG_CARD)
add_text(s5c, "Service Worker (sw.js)",
         Inches(0.55), Inches(2.4), Inches(5.7), Inches(0.45),
         font_size=Pt(15), bold=True, color=RGBColor(0x4C,0xAF,0x50))

sw_steps = [
    ("Install", WHITE,
     "Se activa la primera vez. Pre-cachea index.html y todas las librerías de MediaPipe (~8-15 MB)."),
    ("Activate", GREY_LT,
     "Limpia versiones antiguas del caché para liberar espacio en disco."),
    ("Fetch (cache-first)", CYAN,
     "Cada petición de recurso va primero al caché. Si no está, lo descarga de la red y lo guarda."),
    ("Resultado", RGBColor(0x4C,0xAF,0x50),
     "Después del primer uso, la app arranca y funciona completamente sin conexión a internet."),
]

sw_y = Inches(3.0)
for step_name, color, desc in sw_steps:
    add_rect(s5c, Inches(0.55), sw_y, Inches(1.5), Inches(0.35), fill_color=color)
    add_text(s5c, step_name, Inches(0.55), sw_y, Inches(1.5), Inches(0.35),
             font_size=Pt(11), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s5c, desc, Inches(2.2), sw_y - Pt(2), Inches(4.0), Inches(0.65),
             font_size=Pt(11), color=GREY_LT, wrap=True)
    sw_y += Inches(0.82)

# Right: manifest + ventajas
add_rect(s5c, Inches(6.8), Inches(2.3), Inches(6.1), Inches(2.0), fill_color=BG_CARD)
add_text(s5c, "manifest.json — Instalación",
         Inches(6.95), Inches(2.4), Inches(5.8), Inches(0.45),
         font_size=Pt(15), bold=True, color=RGBColor(0xFF,0xA0,0x00))
add_bullet_box(s5c,
    ['name: "PoseDetection"',
     'display: "fullscreen" — sin barra de navegador',
     'orientation: "portrait-primary"',
     'theme_color: "#0d0d0d" — barra de estado oscura',
     'Icono 192×192 y 512×512 para pantalla inicio'],
    Inches(6.95), Inches(2.9), Inches(5.8), Inches(1.3),
    font_size=Pt(11.5))

add_rect(s5c, Inches(6.8), Inches(4.5), Inches(6.1), Inches(2.2), fill_color=BG_CARD)
add_text(s5c, "Ventajas para Investigación de Campo",
         Inches(6.95), Inches(4.6), Inches(5.8), Inches(0.45),
         font_size=Pt(15), bold=True, color=CYAN)
add_bullet_box(s5c,
    ["Sin dependencia de red en pabellón deportivo",
     "Instalable en móvil del investigador en segundos",
     "Pantalla completa: más espacio para el canvas",
     "Actualización automática al haber nueva versión",
     "Compatible con iOS, Android y escritorio"],
    Inches(6.95), Inches(5.1), Inches(5.8), Inches(1.5),
    font_size=Pt(11.5))

slide_number(s5c, 6, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – APIs del Navegador
# ═══════════════════════════════════════════════════════════════════════════════
s5d = prs.slides.add_slide(blank_layout)
fill_bg(s5d, BG_DARK)
accent_bar(s5d)

add_text(s5d, "APIs del Navegador Utilizadas", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

apis = [
    ("getUserMedia()", CYAN,
     "Acceso a la cámara",
     ["Solicita permiso de cámara al usuario",
      "Configura resolución ideal: 640×480 @ 30fps",
      "Soporta cámara frontal y trasera (facingMode)",
      "Requiere HTTPS o localhost por seguridad",
      "Devuelve un MediaStream con el vídeo en vivo"]),
    ("Canvas 2D API", RGBColor(0xFF,0xA0,0x00),
     "Renderizado en tiempo real",
     ["drawImage() copia cada frame del <video> al canvas",
      "Aplica transformación de espejo (cámara frontal)",
      "Dibuja 35 líneas de esqueleto por frame",
      "Dibuja 33 círculos de landmark con colores",
      "captureStream() convierte canvas a flujo de vídeo"]),
    ("MediaRecorder API", RGBColor(0xFF,0x57,0x22),
     "Codificación de vídeo",
     ["Recibe el stream del canvas en tiempo real",
      "Codifica en MP4 (H.264) o WebM (VP9/VP8)",
      "Genera chunks de datos cada 200 ms",
      "ondataavailable acumula los trozos en un array",
      "onstop: une chunks → Blob → URL de descarga"]),
    ("File System Access API", RGBColor(0x4C,0xAF,0x50),
     "Escritura directa en disco",
     ["showDirectoryPicker(): el usuario elige carpeta",
      "getFileHandle(): abre o crea fichero en esa carpeta",
      "createWritable(): flujo de escritura directo a disco",
      "Solo disponible en Chrome/Edge (Chromium)",
      "Fallback a descarga automática en otros navegadores"]),
]

api_w = Inches(2.9)
api_h = Inches(4.9)
api_gap = Inches(0.35)
api_y = Inches(1.1)
api_x_start = (SLIDE_W - 4*api_w - 3*api_gap) / 2

for i, (name, color, subtitle, items) in enumerate(apis):
    x = api_x_start + i * (api_w + api_gap)
    add_rect(s5d, x, api_y, api_w, api_h, fill_color=BG_CARD)
    # coloured header band
    add_rect(s5d, x, api_y, api_w, Inches(0.55), fill_color=color)
    add_text(s5d, name, x + Inches(0.08), api_y + Pt(2), api_w - Inches(0.1), Inches(0.5),
             font_size=Pt(13), bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    add_text(s5d, subtitle, x + Inches(0.08), api_y + Inches(0.6), api_w - Inches(0.1), Inches(0.38),
             font_size=Pt(11), bold=False, color=color, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s5d, "▸ " + item,
                 x + Inches(0.1), api_y + Inches(1.1) + j * Inches(0.7),
                 api_w - Inches(0.15), Inches(0.68),
                 font_size=Pt(10.5), color=GREY_LT, wrap=True)

# bottom note
add_text(s5d,
         "▸  performance.now()  proporciona timestamps de alta precisión (sub-milisegundo) para sincronizar "
         "cada landmark con su fotograma exacto en el JSON de metadatos.",
         Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.6),
         font_size=Pt(11), color=GREY_MED, wrap=True)

slide_number(s5d, 7, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Detección de Pose (MediaPipe)
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(blank_layout)
fill_bg(s5, BG_DARK)
accent_bar(s5)

add_text(s5, "Detección de Pose con MediaPipe", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Left column: landmarks explanation
add_rect(s5, Inches(0.4), Inches(1.0), Inches(5.8), Inches(5.8), fill_color=BG_CARD)
add_text(s5, "33 Landmarks del Cuerpo Humano",
         Inches(0.55), Inches(1.1), Inches(5.5), Inches(0.5),
         font_size=Pt(15), bold=True, color=RED)

landmark_groups = [
    ("Cara (11 puntos)", CYAN,
     "Ojos, nariz, orejas, boca"),
    ("Tronco superior (8 puntos)", WHITE,
     "Hombros, codos, muñecas"),
    ("Manos (6 puntos/mano)", GREY_LT,
     "Puntas de los dedos principales"),
    ("Tronco inferior (8 puntos)", WHITE,
     "Caderas, rodillas, tobillos"),
    ("Pies (4 puntos/pie)", GREY_LT,
     "Talón, índice del pie"),
]

y_offset = Inches(1.7)
for name, color, detail in landmark_groups:
    add_rect(s5, Inches(0.55), y_offset, Inches(0.12), Inches(0.3), fill_color=color)
    add_text(s5, name, Inches(0.8), y_offset - Pt(2), Inches(5.1), Inches(0.35),
             font_size=Pt(13), bold=True, color=color)
    add_text(s5, detail, Inches(0.8), y_offset + Inches(0.3), Inches(5.1), Inches(0.32),
             font_size=Pt(12), color=GREY_LT)
    y_offset += Inches(0.82)

# landmark data fields
add_text(s5, "Datos por landmark:",
         Inches(0.55), Inches(5.95), Inches(5.5), Inches(0.4),
         font_size=Pt(13), bold=True, color=RED)
add_text(s5, "  x, y  (normalizado 0.0–1.0)  ·  z  (profundidad relativa)  ·  visibility  (0.0–1.0)",
         Inches(0.55), Inches(6.35), Inches(5.5), Inches(0.4),
         font_size=Pt(12), color=GREY_LT)

# Right column: rendering
add_rect(s5, Inches(6.6), Inches(1.0), Inches(6.3), Inches(5.8), fill_color=BG_CARD)
add_text(s5, "Renderizado sobre Canvas",
         Inches(6.75), Inches(1.1), Inches(6.0), Inches(0.5),
         font_size=Pt(15), bold=True, color=RED)

render_items = [
    ("Esqueleto", WHITE,
     "35 conexiones · líneas blancas semi-transparentes · 2.5 px"),
    ("Landmarks cara", CYAN,
     "Círculos cian · radio 4-5 px · visibilidad ≥ 0.3"),
    ("Landmarks cuerpo", RGBColor(0xFF,0x57,0x22),
     "Círculos rojo/naranja · radio 4 px"),
    ("Espejado", GREY_LT,
     "Solo en render (cámara frontal) · JSON sin espejo"),
    ("Frecuencia", RGBColor(0xFF,0xA0,0x00),
     "requestAnimationFrame · ~30 fps objetivo"),
]

ry = Inches(1.7)
for rname, rcolor, rdesc in render_items:
    add_rect(s5, Inches(6.75), ry, Inches(0.12), Inches(0.3), fill_color=rcolor)
    add_text(s5, rname, Inches(7.0), ry - Pt(2), Inches(5.7), Inches(0.35),
             font_size=Pt(13), bold=True, color=rcolor)
    add_text(s5, rdesc, Inches(7.0), ry + Inches(0.3), Inches(5.7), Inches(0.32),
             font_size=Pt(12), color=GREY_LT)
    ry += Inches(0.9)

slide_number(s5, 8, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – Formato de Salida: JSON
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(blank_layout)
fill_bg(s6, BG_DARK)
accent_bar(s6)

add_text(s6, "Formato de Salida — Metadatos JSON", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# JSON schema box
add_rect(s6, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.8), fill_color=BG_CARD)
add_text(s6, "Estructura del fichero JSON",
         Inches(0.55), Inches(1.1), Inches(5.7), Inches(0.45),
         font_size=Pt(14), bold=True, color=RED)

json_code = """{
  "recording": {
    "id": "rec_20260312_143022_ab3x",
    "dorsal": "07",
    "posture_code": 1,
    "posture_name": "tiro_libre",
    "date": "2026-03-12T14:30:22.000Z",
    "duration_seconds": 12.34,
    "fps": 28.5,
    "resolution": { "width": 640, "height": 480 },
    "platform": "mobile",
    "video_mime": "video/mp4;codecs=avc1"
  },
  "frames": [
    {
      "frame_index": 0,
      "timestamp_ms": 0.00,
      "landmarks": [
        {
          "index": 0,
          "name": "nose",
          "x": 0.512345,
          "y": 0.318712,
          "z": -0.087234,
          "visibility": 0.998100
        }
        // ... 32 landmarks más
      ]
    }
    // ... N fotogramas
  ]
}"""

add_text(s6, json_code,
         Inches(0.55), Inches(1.65), Inches(5.7), Inches(5.0),
         font_size=Pt(9.5), color=RGBColor(0xA5,0xD6,0xA7),
         wrap=True)

# Right column: key facts
add_rect(s6, Inches(6.8), Inches(1.0), Inches(6.1), Inches(2.6), fill_color=BG_CARD)
add_bullet_box(s6,
    ["ID único: timestamp + sufijo aleatorio",
     "Fecha ISO 8601 + duración real",
     "FPS promedio medido durante grabación",
     "User-Agent para identificar dispositivo",
     "Resolución: 640 × 480 px siempre"],
    Inches(6.95), Inches(1.05), Inches(5.8), Inches(2.5),
    title="Bloque recording", font_size=Pt(12))

add_rect(s6, Inches(6.8), Inches(3.8), Inches(6.1), Inches(3.0), fill_color=BG_CARD)
add_bullet_box(s6,
    ["Un objeto por fotograma capturado",
     "timestamp_ms con performance.now()",
     "33 landmarks × fotograma",
     "6 decimales de precisión (x, y, z, vis)",
     "visibility en JSON sin umbralizar",
     "Tamaño típico: 500 KB – 2 MB"],
    Inches(6.95), Inches(3.85), Inches(5.8), Inches(2.9),
    title="Array frames", font_size=Pt(12))

slide_number(s6, 9, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 – Interfaz de Usuario
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(blank_layout)
fill_bg(s7, BG_DARK)
accent_bar(s7)

add_text(s7, "Interfaz de Usuario", Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

ui_sections = [
    ("Panel de Control", RED, [
        "Selector de dorsal: 00-99",
        "Código de postura: 7 tipos de movimiento",
        "Toggle auto-guardar",
        "Selector de carpeta destino (Chromium)",
    ]),
    ("Zona de Grabación", CYAN, [
        "Canvas 640×480 con overlay de esqueleto",
        "Botón REC con indicador de tiempo",
        "Límite visual de 30 segundos",
        "Botón flip cámara (frontal/trasera)",
    ]),
    ("Panel de Descarga", RGBColor(0xFF,0xA0,0x00), [
        "Aparece al terminar la grabación",
        "Botón descargar vídeo (.mp4/.webm)",
        "Botón descargar metadatos (.json)",
        "Fallback manual para iOS/Firefox",
    ]),
    ("Diseño Visual", RGBColor(0x4C,0xAF,0x50), [
        "Tema oscuro: fondo #0d0d0d",
        "Acento rojo #e53935",
        "Mobile-first, sin delay 300ms",
        "Safe area insets (notch/barra inicio)",
    ]),
]

card_w = Inches(5.7)
card_h = Inches(2.6)
positions_ui = [
    (Inches(0.4),  Inches(1.05)),
    (Inches(7.2),  Inches(1.05)),
    (Inches(0.4),  Inches(3.85)),
    (Inches(7.2),  Inches(3.85)),
]

for (x, y), (title, color, items) in zip(positions_ui, ui_sections):
    add_rect(s7, x, y, card_w, card_h, fill_color=BG_CARD)
    add_bullet_box(s7, items, x + Inches(0.12), y + Inches(0.05),
                   card_w - Inches(0.2), card_h - Inches(0.1),
                   title=title, title_color=color, font_size=Pt(12.5))

slide_number(s7, 10, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 – Compatibilidad y Rendimiento
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(blank_layout)
fill_bg(s8, BG_DARK)
accent_bar(s8)

add_text(s8, "Compatibilidad y Rendimiento", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

# Compatibility table header
add_rect(s8, Inches(0.4), Inches(1.0), Inches(5.9), Inches(5.6), fill_color=BG_CARD)
add_text(s8, "Compatibilidad de Navegadores",
         Inches(0.55), Inches(1.1), Inches(5.7), Inches(0.45),
         font_size=Pt(14), bold=True, color=RED)

browsers = [
    ("Chrome / Edge 90+",   "✓ Auto-save · ✓ MP4",       RGBColor(0x4C,0xAF,0x50)),
    ("Firefox 90+",          "✓ Descarga manual · WebM",  RGBColor(0xFF,0xA0,0x00)),
    ("Safari 15+ (iOS/Mac)", "✓ Descarga manual · MP4",   RGBColor(0xFF,0xA0,0x00)),
    ("Requiere HTTPS",        "o localhost para cámara",   GREY_MED),
]

by = Inches(1.7)
for browser, notes, color in browsers:
    add_rect(s8, Inches(0.55), by, Inches(0.14), Inches(0.35), fill_color=color)
    add_text(s8, browser, Inches(0.82), by - Pt(2), Inches(2.6), Inches(0.45),
             font_size=Pt(13), bold=True, color=color)
    add_text(s8, notes, Inches(3.55), by - Pt(2), Inches(2.5), Inches(0.45),
             font_size=Pt(12), color=GREY_LT)
    by += Inches(0.72)

# Performance / limits
add_rect(s8, Inches(0.4), Inches(4.0), Inches(5.9), Inches(2.6), fill_color=BG_CARD)
add_bullet_box(s8,
    ["Resolución fija: 640 × 480 px",
     "FPS objetivo: ~30 (varía según dispositivo)",
     "Grabación máxima: 30 segundos",
     "Tamaño vídeo: 2-8 MB típico",
     "Tamaño JSON: 500 KB – 2 MB"],
    Inches(0.55), Inches(4.05), Inches(5.7), Inches(2.5),
    title="Parámetros de Rendimiento", font_size=Pt(12))

# Right: known limitations & offline
add_rect(s8, Inches(6.7), Inches(1.0), Inches(6.2), Inches(2.8), fill_color=BG_CARD)
add_bullet_box(s8,
    ["Service Worker — cache-first",
     "Pre-cache de librerías MediaPipe (8-15 MB)",
     "Primer uso requiere internet",
     "Usos posteriores: 100% offline",
     "PWA instalable como app nativa"],
    Inches(6.85), Inches(1.05), Inches(6.0), Inches(2.7),
    title="Soporte Offline (PWA)", title_color=RGBColor(0x4C,0xAF,0x50), font_size=Pt(12))

add_rect(s8, Inches(6.7), Inches(4.0), Inches(6.2), Inches(2.6), fill_color=BG_CARD)
add_bullet_box(s8,
    ["Sin gestión de caída de frames",
     "Resolución no adaptativa",
     "Sin audio en la grabación",
     "iOS: reinicio manual si se bloquea pantalla",
     "File System API solo en Chromium"],
    Inches(6.85), Inches(4.05), Inches(6.0), Inches(2.5),
    title="Limitaciones Conocidas", title_color=RGBColor(0xFF,0x57,0x22), font_size=Pt(12))

slide_number(s8, 11, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 – Flujo de Uso
# ═══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(blank_layout)
fill_bg(s9, BG_DARK)
accent_bar(s9)

add_text(s9, "Flujo de Uso", Inches(0.5), Inches(0.18), Inches(10), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

steps = [
    ("1", "Preparación",
     ["Abrir app en Chrome/Edge (HTTPS)",
      "Conceder permiso de cámara",
      "Seleccionar carpeta de guardado (opcional)"]),
    ("2", "Configuración",
     ["Elegir dorsal del jugador (00-99)",
      "Seleccionar tipo de movimiento (1-7)",
      "Activar auto-guardado si se desea"]),
    ("3", "Grabación",
     ["Pulsar botón REC",
      "El esqueleto se dibuja en tiempo real",
      "Máximo 30 s — se para automáticamente"]),
    ("4", "Exportación",
     ["Fichero vídeo: {dorsal}_{ts}_{postura}.mp4",
      "Fichero JSON: misma base + _metadata",
      "Auto-guardar en carpeta o descarga manual"]),
    ("5", "Análisis",
     ["JSON contiene todos los fotogramas",
      "33 landmarks × frame → series temporales",
      "Procesado post en Python/R/MATLAB"]),
]

step_w = Inches(2.35)
step_h = Inches(5.0)
step_gap = Inches(0.25)
sx = (SLIDE_W - (len(steps)*step_w + (len(steps)-1)*step_gap)) / 2
sy = Inches(1.1)

for i, (num, title, items) in enumerate(steps):
    x = sx + i * (step_w + step_gap)
    add_rect(s9, x, sy, step_w, step_h, fill_color=BG_CARD)
    # number circle
    add_rect(s9, x + Inches(0.85), sy + Inches(0.1), Inches(0.6), Inches(0.6), fill_color=RED)
    add_text(s9, num, x + Inches(0.85), sy + Inches(0.1), Inches(0.6), Inches(0.6),
             font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s9, title, x + Inches(0.1), sy + Inches(0.85), step_w - Inches(0.2), Inches(0.45),
             font_size=Pt(13), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(s9, "• " + item,
                 x + Inches(0.12), sy + Inches(1.45) + j*Inches(0.85),
                 step_w - Inches(0.2), Inches(0.82),
                 font_size=Pt(10.5), color=GREY_LT, wrap=True)

slide_number(s9, 12, 13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 – Conclusiones y Trabajo Futuro
# ═══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(blank_layout)
fill_bg(s10, BG_DARK)
accent_bar(s10)

add_text(s10, "Conclusiones y Trabajo Futuro", Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
         font_size=Pt(28), bold=True, color=WHITE)

add_rect(s10, Inches(0.4), Inches(1.0), Inches(6.0), Inches(5.6), fill_color=BG_CARD)
add_bullet_box(s10,
    ["Herramienta 100% cliente — sin infraestructura de servidor",
     "Funciona en cualquier smartphone moderno",
     "Captura datos biomecánicos con precisión de 6 dígitos",
     "Soporte offline via PWA y Service Worker",
     "Formato JSON estandarizado para post-proceso",
     "Integración directa con flujos de investigación",
     "Etiquetado especializado para baloncesto (7 gestos)",
     "Despliegue trivial: un solo fichero HTML"],
    Inches(0.55), Inches(1.05), Inches(5.8), Inches(5.5),
    title="Logros Clave", font_size=Pt(12.5))

add_rect(s10, Inches(6.8), Inches(1.0), Inches(6.1), Inches(5.6), fill_color=BG_CARD)
add_bullet_box(s10,
    ["Análisis multi-persona (actualmente 1 persona)",
     "Detección de fases de movimiento automática",
     "Integración con pipeline ML (TensorFlow.js)",
     "Dashboard de análisis biomecánico en tiempo real",
     "Exportación a formatos BVH / C3D estándar",
     "Soporte de audio y etiquetado vocal",
     "Comparativa entre jugadores / sesiones",
     "API REST para sincronización con servidor"],
    Inches(6.95), Inches(1.05), Inches(5.8), Inches(5.5),
    title="Trabajo Futuro", title_color=CYAN, font_size=Pt(12.5))

slide_number(s10, 13, 13)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = "/home/jordi/Documentos/DoctoradoUVa/Development/poseDetection/PoseDetection_Presentacion.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
